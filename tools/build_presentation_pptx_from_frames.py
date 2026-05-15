#!/usr/bin/env python3

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path.cwd()
FRAMES_DIR = ROOT / "exports" / "presentation" / "frames"
OUT = ROOT / "exports" / "presentation" / "joseph-robinson-ist782-portfolio-presentation.pptx"

ROOT_URL = "https://blackbone.live"
PORTFOLIO_URL = "https://blackbone.live/portfolio/"

SLIDE_URLS = {
    1: ROOT_URL,
    2: PORTFOLIO_URL,
    3: "https://blackbone.live/portfolio/reader/?doc=financial-packet",
    4: "https://blackbone.live/portfolio/reader/?doc=chemical-packet",
    5: "https://blackbone.live/portfolio/reader/?doc=inventory-packet",
    6: PORTFOLIO_URL,
    7: PORTFOLIO_URL,
    8: PORTFOLIO_URL,
    9: PORTFOLIO_URL,
}

# Approximate clickable rectangles over visible link text.
# Full-slide image is also linked, so these are extra reinforcement.
LINK_RECTS = {
    1: (3.10, 7.05, 3.10, 0.45),  # https://blackbone.live under degree line
    3: (3.00, 5.55, 2.25, 0.45),  # View Evidence
    4: (3.00, 5.55, 2.25, 0.45),  # View Evidence
    5: (3.00, 6.20, 2.25, 0.45),  # View Evidence
}

frames = sorted(FRAMES_DIR.glob("slide-*.png"))

if not frames:
    raise SystemExit(f"No frames found in {FRAMES_DIR}")

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_layout = prs.slide_layouts[6]

def add_click_rect(slide, url, rect):
    x, y, w, h = rect
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    box.click_action.hyperlink.address = url
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 0, 0)

    try:
        box.fill.transparency = 100
    except Exception:
        pass

    try:
        box.line.fill.background()
    except Exception:
        pass

for frame in frames:
    slide_num = int(frame.stem.split("-")[-1])
    url = SLIDE_URLS.get(slide_num, PORTFOLIO_URL)

    print(f"Adding {frame.name} -> {url}")

    slide = prs.slides.add_slide(blank_layout)

    pic = slide.shapes.add_picture(
        str(frame),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

    pic.click_action.hyperlink.address = url

    if slide_num in LINK_RECTS:
        add_click_rect(slide, url, LINK_RECTS[slide_num])

prs.save(OUT)

print()
print(f"Created: {OUT}")
print("Clickable behavior: test in slideshow mode, or Ctrl-click in edit mode depending on viewer.")
